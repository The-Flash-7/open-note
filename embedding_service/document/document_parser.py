# Copyright (c) 2026 litongshuai
# SPDX-License-Identifier: MIT OR Apache-2.0

import os
from typing import Optional, Tuple


class DocumentParser:
    SUPPORTED_EXTENSIONS = ['.pdf', '.docx', '.pptx']
    
    @staticmethod
    def get_file_extension(file_path: str) -> str:
        _, ext = os.path.splitext(file_path)
        return ext.lower()
    
    @staticmethod
    def is_supported(file_path: str) -> bool:
        ext = DocumentParser.get_file_extension(file_path)
        return ext in DocumentParser.SUPPORTED_EXTENSIONS
    
    @staticmethod
    def parse_pdf(file_path: str) -> Tuple[Optional[str], Optional[str]]:
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(file_path)
            text_content = []
            
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
            
            full_text = '\n\n'.join(text_content)
            
            metadata = {}
            if reader.metadata:
                metadata = {
                    'title': reader.metadata.get('/Title', ''),
                    'author': reader.metadata.get('/Author', ''),
                    'creator': reader.metadata.get('/Creator', ''),
                    'producer': reader.metadata.get('/Producer', ''),
                    'creation_date': str(reader.metadata.get('/CreationDate', '')),
                    'modification_date': str(reader.metadata.get('/ModDate', '')),
                }
            
            return full_text.strip(), metadata
            
        except Exception as e:
            return None, f"PDF parsing error: {str(e)}"
    
    @staticmethod
    def parse_docx(file_path: str) -> Tuple[Optional[str], Optional[str]]:
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(' | '.join(row_text))
            
            full_text = '\n\n'.join(text_content)
            
            core_props = doc.core_properties
            metadata = {
                'title': getattr(core_props, 'title', None) or '',
                'author': getattr(core_props, 'author', None) or '',
                'subject': getattr(core_props, 'subject', None) or '',
                'keywords': getattr(core_props, 'keywords', None) or '',
                'last_modified_by': getattr(core_props, 'last_modified_by', None) or '',
                'created': str(getattr(core_props, 'created', None) or ''),
                'modified': str(getattr(core_props, 'modified', None) or ''),
            }
            
            return full_text.strip(), metadata
            
        except Exception as e:
            return None, f"Word parsing error: {str(e)}"
    
    @staticmethod
    def parse_pptx(file_path: str) -> Tuple[Optional[str], Optional[str]]:
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    text_content.append(f"--- Slide {slide_num} ---\n" + '\n\n'.join(slide_texts))
            
            full_text = '\n\n'.join(text_content)
            
            core_props = prs.core_properties
            metadata = {
                'title': getattr(core_props, 'title', None) or '',
                'author': getattr(core_props, 'author', None) or '',
                'subject': getattr(core_props, 'subject', None) or '',
                'keywords': getattr(core_props, 'keywords', None) or '',
                'last_modified_by': getattr(core_props, 'last_modified_by', None) or '',
                'created': str(getattr(core_props, 'created', None) or ''),
                'modified': str(getattr(core_props, 'modified', None) or ''),
                'slide_count': len(prs.slides),
            }
            
            return full_text.strip(), metadata
            
        except Exception as e:
            return None, f"PPT parsing error: {str(e)}"
    
    @staticmethod
    def parse_document(file_path: str) -> dict:
        ext = DocumentParser.get_file_extension(file_path)
        
        if not os.path.exists(file_path):
            return {
                'success': False,
                'error': f"File not found: {file_path}",
                'text': None,
                'metadata': None
            }
        
        if not DocumentParser.is_supported(file_path):
            return {
                'success': False,
                'error': f"Unsupported file type: {ext}",
                'text': None,
                'metadata': None
            }
        
        text = None
        metadata = None
        error = None
        
        if ext == '.pdf':
            text, result = DocumentParser.parse_pdf(file_path)
            if text is None:
                error = result
            else:
                metadata = result
        
        elif ext == '.docx':
            text, result = DocumentParser.parse_docx(file_path)
            if text is None:
                error = result
            else:
                metadata = result
        
        elif ext == '.pptx':
            text, result = DocumentParser.parse_pptx(file_path)
            if text is None:
                error = result
            else:
                metadata = result
        
        if error:
            return {
                'success': False,
                'error': error,
                'text': None,
                'metadata': None
            }
        
        return {
            'success': True,
            'error': None,
            'text': text,
            'metadata': metadata,
            'file_type': ext
        }